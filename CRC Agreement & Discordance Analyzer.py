import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
import itertools, json, os
import statsmodels.api as sm
from statsmodels.stats.contingency_tables import mcnemar as sm_mcnemar
from sklearn.metrics import roc_auc_score, roc_curve
from scipy import stats
from zipfile import ZipFile, ZIP_DEFLATED

# --- VIF: ê²¬ê³ í•œ ì„í¬íŠ¸(ì—†ìœ¼ë©´ ê²½ê³ ë§Œ) ---
try:
    from statsmodels.stats.outliers_influence import variance_inflation_factor  # <- ì •í™•í•œ ê²½ë¡œ
    HAS_VIF = True
except Exception:
    variance_inflation_factor = None
    HAS_VIF = False

# optional deps
try:
    import plotly.express as px
    import plotly.graph_objects as go
    HAS_PLOTLY = True
except Exception:
    HAS_PLOTLY = False
try:
    import pingouin as pg
    HAS_PINGOUIN = True
except Exception:
    HAS_PINGOUIN = False

# ---------------- App meta ----------------
st.set_page_config(page_title="CRC Agreement Analyzer â€” MAX", layout="wide")
st.title("CRC Agreement Analyzer â€” MAX (Everything-in-One)")

# ---------------- Session-safe stores ----------------
for k, default in [
    ("results_tables", {}),
    ("consensus_map", {}),
    ("images_store", {}),
]:
    if k not in st.session_state:
        st.session_state[k] = default

# ---------------- Glossary (ê°„ë‹¨ ì •ì˜) ----------------
GLOSSARY = {
    "Light's Îº (weighted mean)": "ì—¬ëŸ¬ í‰ê°€ì ìŒì˜ ê°€ì¤‘ Cohen Îº í‰ê· (ê°€ì¤‘ì²´ê³„ ì„ íƒ).",
    "Fleiss' Îº (unweighted)": "ë‹¤ìˆ˜ í‰ê°€ìì˜ ë²”ì£¼ ì¼ì¹˜ë„(ëª¨ë“  ìŒì„ í•œêº¼ë²ˆì—).",
    "Krippendorff's Î± (ordinal)": "ì„œì—´/ëª…ëª©/ê°„ê²© ë“±ê¸‰ì— ë§ì¶˜ ì¼ë°˜í™” ì¼ì¹˜ë„.",
    "ICC(2,k)": "í‰ê°€ì ê°„ ì¼ì¹˜ë„(í˜¼í•©íš¨ê³¼Â·ì ˆëŒ€í•©ì¹˜Â·í‰ê· ì¸¡ì •).",
    "Brennanâ€“Prediger Îº": "ìš°ë„ì— ëœ ë¯¼ê°í•œ Îº(ê¸°ëŒ€í•©ì¹˜=1/k).",
    "PABAK (multi)": "ê´€ì°°í•©ì¹˜(Po) ê¸°ë°˜ì˜ ë³´ì • Îº(=2Â·Poâˆ’1).",
    "McNemar": "ë‘ ì´ì§„ ë¶„ë¥˜ì˜ ë¶ˆì¼ì¹˜(01 vs 10) ë¹„êµ ê²€ì •.",
    "AUC": "ROC ê³¡ì„  ì•„ë˜ ë©´ì (ì„ê³„ì¹˜ ë¶ˆë³€ì˜ ìˆœìœ„ ê¸°ë°˜ ì„±ëŠ¥).",
    "DeLong": "ë‘ AUC ì°¨ì´ì˜ ë¹„ëª¨ìˆ˜ ì¶”ì •/ê²€ì •.",
    "Firth logistic": "í¬ê·€ì‚¬ê±´/ì™„ì „ë¶„ë¦¬ ë³´ì • ë¡œì§€ìŠ¤í‹± íšŒê·€.",
    "VIF": "ê³µì„ ì„± ì§€í‘œ(>5~10ì´ë©´ ë‹¤ì¤‘ê³µì„ ì„± ì˜ì‹¬).",
}

# ---------------- Utilities ----------------
@st.cache_data
def read_excel(file):
    try:
        return pd.read_excel(file)
    except Exception as e:
        st.error(f"ì—‘ì…€ íŒŒì¼ ì½ê¸° ì˜¤ë¥˜: {e}")
        return None

def safe_num(s): return pd.to_numeric(s, errors="coerce")

def interpret_kappa(k):
    if pd.isna(k): return "N/A"
    return "Poor" if k<0 else ("Slight" if k<=0.20 else ("Fair" if k<=0.40 else ("Moderate" if k<=0.60 else ("Substantial" if k<=0.80 else "Almost Perfect"))))

def make_distance_matrix(k, scheme="quadratic"):
    D = np.zeros((k,k), dtype=float)
    for i in range(k):
        for j in range(k):
            d = abs(i-j)/(k-1) if k>1 else 0.0
            if scheme=="linear": D[i,j]=d
            elif scheme=="quadratic": D[i,j]=d**2
            elif scheme=="stepwise":
                D[i,j] = 0.0 if d==0 else (0.33 if d<=0.25 else (0.66 if d<=0.5 else 1.0))
            elif scheme=="unweighted": D[i,j] = 0.0 if i==j else 1.0
            else: D[i,j]=d**2
    return D

def sanitize_custom_D(W):
    W = np.array(W, dtype=float)
    if W.shape[0]!=W.shape[1]: raise ValueError("ì»¤ìŠ¤í…€ ê°€ì¤‘í–‰ë ¬ì€ ì •ì‚¬ê°(kÃ—k)ì´ì–´ì•¼ í•©ë‹ˆë‹¤.")
    W = (W + W.T)/2.0
    np.fill_diagonal(W, 0.0)
    W = np.clip(W, 0.0, 1.0)
    return W

def weighted_kappa_custom(a, b, labels, D):
    df = pd.DataFrame({"a":a, "b":b}).dropna()
    if df.empty: return np.nan
    lab_to_i = {lab:i for i,lab in enumerate(labels)}
    ai = df["a"].map(lab_to_i); bi=df["b"].map(lab_to_i)
    k = len(labels)
    cm = np.zeros((k,k), dtype=float)
    for i_, j_ in zip(ai,bi):
        if pd.notna(i_) and pd.notna(j_): cm[int(i_), int(j_)] += 1
    n = cm.sum()
    if n==0: return np.nan
    r = cm.sum(axis=1); c = cm.sum(axis=0)
    Do = (D*cm).sum()
    E = np.outer(r,c)/n
    De = (D*E).sum()
    return float(1.0 - Do/De) if De!=0 else np.nan

def bootstrap_ci_stat(fun, n, B=2000, seed=42):
    rng = np.random.default_rng(seed)
    vals = []
    for _ in range(B):
        idx = rng.integers(0, n, n)
        vals.append(fun(idx))
    vals = np.array(vals, dtype=float)
    lo, hi = np.nanpercentile(vals, [2.5,97.5])
    p = 2*min((vals<=0).mean(), (vals>=0).mean())
    return float(lo), float(hi), float(p), vals

def bootstrap_multi_stat(df, stat_func, B=2000, seed=42):
    n = len(df)
    if n == 0: return np.nan, np.nan, np.array([])
    rng = np.random.default_rng(seed)
    vals = [stat_func(df.iloc[rng.integers(0, n, n)]) for _ in range(B)]
    vals = np.array(vals, dtype=float)
    lo, hi = np.nanpercentile(vals, [2.5, 97.5])
    return float(lo), float(hi), vals

def fleiss_kappa(ratings_df, B=2000, seed=42):
    X=ratings_df.copy()
    cats = sorted([c for c in pd.unique(X.values.ravel()) if pd.notna(c)])
    if len(cats)==0: return np.nan, np.nan, np.nan, np.array([])
    m=len(cats); cat_to_i={c:i for i,c in enumerate(cats)}
    n_items = X.shape[0]
    def calc(df_sub):
        n_items_sub = df_sub.shape[0]
        if n_items_sub==0: return np.nan
        N = np.zeros((n_items_sub, m), dtype=float)
        for u in range(n_items_sub):
            vc = pd.Series(df_sub.iloc[u,:]).dropna().value_counts()
            for c, cnt in vc.items():
                if c in cat_to_i: N[u,cat_to_i[c]]=cnt
        n_i = N.sum(axis=1)
        mask = n_i >= 2
        if not np.any(mask): return np.nan
        N = N[mask,:]; n_i = n_i[mask]
        Pu = ( (N*(N-1)).sum(axis=1) ) / (n_i*(n_i-1))
        P_bar = np.nanmean(Pu)
        total = n_i.sum()
        if total==0: return np.nan
        pj = N.sum(axis=0)/total
        Pe = (pj**2).sum()
        return (P_bar-Pe)/(1-Pe) if (1-Pe)!=0 else np.nan
    k0 = calc(X)
    rng=np.random.default_rng(seed)
    boots=[calc(X.iloc[rng.integers(0,n_items,n_items)]) for _ in range(B)]
    lo, hi = np.nanpercentile(np.array(boots,float), [2.5,97.5])
    return float(k0), float(lo), float(hi), np.array(boots)

def krippendorff_alpha(ratings_df, level="ordinal", B=2000, seed=42):
    X=ratings_df.copy()
    cats = sorted([c for c in pd.unique(X.values.ravel()) if pd.notna(c)])
    if len(cats)==0: return np.nan, np.nan, np.nan, np.array([])
    m=len(cats); cat_to_i={c:i for i,c in enumerate(cats)}
    def calc(df):
        O = np.zeros((m,m), dtype=float)
        for _, row in df.iterrows():
            vals = [v for v in row.values if pd.notna(v)]
            if len(vals)<2: continue
            vc = pd.Series(vals).value_counts()
            for i_c, ci in vc.items():
                ii=cat_to_i[i_c]; O[ii,ii] += ci*(ci-1)
                for j_c, cj in vc.items():
                    if j_c!=i_c: O[ii,cat_to_i[j_c]] += ci*cj
        if O.sum()==0: return np.nan
        D = np.zeros((m,m), dtype=float)
        for i in range(m):
            for j in range(m):
                if level=="nominal": D[i,j] = 0.0 if i==j else 1.0
                else:
                    d=abs(i-j)/(m-1) if m>1 else 0.0; D[i,j]=d**2
        Do = (O*D).sum()/O.sum()
        n_i = O.sum(axis=1); E = np.outer(n_i,n_i)
        for i in range(m): E[i,i]=n_i[i]*(n_i[i]-1)
        if E.sum()==0: return np.nan
        De = (E*D).sum()/E.sum()
        return 1.0 - Do/De if De!=0 else np.nan
    a0 = calc(X)
    rng=np.random.default_rng(seed)
    boots=[calc(X.iloc[rng.integers(0,len(X),len(X))]) for _ in range(B)]
    lo, hi = np.nanpercentile(np.array(boots,float), [2.5,97.5])
    return float(a0), float(lo), float(hi), np.array(boots)

def observed_exact_agreement_multi(raters_df):
    agree_pairs=0.0; total_pairs=0.0
    for _, row in raters_df.iterrows():
        vals = pd.Series(row).dropna().values
        n = len(vals)
        if n<2: continue
        total_pairs += n*(n-1)/2.0
        vc = pd.Series(vals).value_counts()
        agree_pairs += sum(c*(c-1)/2.0 for c in vc.values)
    return (agree_pairs/total_pairs) if total_pairs>0 else np.nan

def brennan_prediger_kappa(raters_df):
    Po = observed_exact_agreement_multi(raters_df)
    cats = sorted([c for c in pd.unique(raters_df.values.ravel()) if pd.notna(c)])
    k = max(1, len(cats))
    return ((Po - 1.0/k) / (1.0 - 1.0/k)) if k>1 and pd.notna(Po) else np.nan

def pabak_multi(raters_df):
    Po = observed_exact_agreement_multi(raters_df)
    return 2*Po - 1 if pd.notna(Po) else np.nan

def pairwise_kappa_table(raters_df, labels, D, B=2000, seed=42):
    rows=[]; n=len(raters_df)
    pvals=[]
    for a,b in itertools.combinations(raters_df.columns, 2):
        stat=lambda idx: weighted_kappa_custom(raters_df[a].iloc[idx], raters_df[b].iloc[idx], labels, D)
        k0 = stat(np.arange(n))
        lo,hi,p,_ = bootstrap_ci_stat(stat, n, B=B, seed=seed)
        pvals.append(p)
        pea = float((raters_df[a]==raters_df[b]).mean())*100.0
        rows.append({"Rater A":a,"Rater B":b,"Exact Agreement (%)":round(pea,2),
                     "Weighted Kappa":round(k0,4),"95% CI Lower":round(lo,4),
                     "95% CI Upper":round(hi,4),"P (bootstrap)":round(p,6),
                     "Interp": interpret_kappa(k0)})
    # Holm ë³´ì •
    m=len(pvals)
    if m>0:
        order=np.argsort(pvals)
        adj=[None]*m; maxv=0
        for rank, idx in enumerate(order):
            padj = min(1.0, pvals[idx]*(m-rank))
            maxv = max(maxv, padj)
            adj[idx]=maxv
        for i,row in enumerate(rows): row["P Holm"] = round(adj[i],6)
    return pd.DataFrame(rows)

def lights_kappa(raters_df, labels, D, B=2000, seed=42):
    pairs=list(itertools.combinations(raters_df.columns,2)); n=len(raters_df)
    if len(pairs)==0: return np.nan, np.nan, np.nan, np.array([])
    def stat(idx):
        vals=[weighted_kappa_custom(raters_df[a].iloc[idx], raters_df[b].iloc[idx], labels, D) for a,b in pairs]
        return float(np.nanmean(vals))
    k0 = stat(np.arange(n))
    lo,hi,p,boots = bootstrap_ci_stat(stat, n, B=B, seed=seed)
    return float(k0), float(lo), float(hi), boots

def calculate_icc(raters_df):
    if not HAS_PINGOUIN:
        return np.nan, np.nan, np.nan
    df_long = raters_df.reset_index().melt(id_vars='index', var_name='rater', value_name='rating').rename(columns={'index':'case'}).dropna()
    if df_long.empty or df_long['case'].nunique()<2 or df_long['rater'].nunique()<2:
        return np.nan, np.nan, np.nan
    try:
        tbl = pg.intraclass_corr(data=df_long, targets='case', raters='rater', ratings='rating').set_index('Type')
        val, (lo,hi) = tbl.loc['ICC2k','ICC'], tbl.loc['ICC2k','CI95%']
        return float(val), float(lo), float(hi)
    except Exception:
        return np.nan, np.nan, np.nan

def accuracy_metric(y_true, y_pred_bin):
    v = ~(y_true.isna() | y_pred_bin.isna())
    return float((y_true[v]==y_pred_bin[v]).mean()) if v.sum()>0 else np.nan

def bootstrap_metric_ci(y_true, y_score, metric_fun, B=2000, seed=42):
    rng = np.random.default_rng(seed); n=len(y_true); vals=[]
    for _ in range(B):
        idx=rng.integers(0,n,n); vals.append(metric_fun(y_true.iloc[idx], y_score.iloc[idx]))
    lo,hi = np.nanpercentile(vals, [2.5,97.5])
    return float(lo), float(hi), np.array(vals)

# --- DeLong helpers ---
def _compute_midrank(x):
    J = np.argsort(x); Z = x[J]; N = len(x); T = np.zeros(N, dtype=float); i=0
    while i<N:
        j=i
        while j<N and Z[j]==Z[i]: j+=1
        T[i:j] = 0.5*(i+j-1)+1; i=j
    T2 = np.empty(N, dtype=float); T2[J]=T; return T2

def _fastDeLong(predictions_sorted_transposed, label_1_count):
    m = label_1_count; n = predictions_sorted_transposed.shape[1]-m
    pos = predictions_sorted_transposed[:, :m]; neg = predictions_sorted_transposed[:, m:]; k = predictions_sorted_transposed.shape[0]
    tx = np.empty([k, m], float); ty = np.empty([k, n], float)
    for r in range(k):
        tx[r,:] = _compute_midrank(pos[r,:]); ty[r,:]=_compute_midrank(neg[r,:])
    tz = np.empty([k, m+n], float)
    for r in range(k): tz[r,:] = _compute_midrank(predictions_sorted_transposed[r,:])
    aucs = tz[:,:m].sum(axis=1)/(m*n) - (m+1.0)/(2.0*n)
    v01 = (tz[:,:m] - tx)/n; v10 = 1.0 - (tz[:,m:] - ty)/m
    sx = np.cov(v01); sy = np.cov(v10); delongcov = sx/m + sy/n
    return aucs, delongcov

def delong_roc_variance(y_true, y_scores):
    order = np.argsort(-y_scores); y_scores=y_scores[order]; y_true=y_true[order]
    label_1_count = int(np.sum(y_true))
    aucs, delongcov = _fastDeLong(y_scores[np.newaxis,:], label_1_count)
    return aucs[0], delongcov

def delong_compare(y_true, s1, s2):
    y_true = np.asarray(y_true, dtype=int); s1=np.asarray(s1, float); s2=np.asarray(s2, float)
    auc1, var1 = delong_roc_variance(y_true, s1); auc2, var2 = delong_roc_variance(y_true, s2)
    cov12 = 0.0
    se = np.sqrt(var1 + var2 - 2*cov12).item()
    z = (auc1 - auc2) / se if se>0 else np.nan
    p = 2*(1 - stats.norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    return float(auc1), float(auc2), float(z), float(p)

# --- Firth logistic (bias-reduced) ---
def firth_logit(X, y, max_iter=200, tol=1e-8):
    X = np.asarray(X, float); y = np.asarray(y, int)
    n,p = X.shape
    beta = np.zeros(p)
    for it in range(max_iter):
        eta = X @ beta; p_hat = 1/(1+np.exp(-eta))
        W = np.clip(p_hat*(1-p_hat), 1e-9, None)
        XWX = X.T * W @ X
        try: XWX_inv = np.linalg.inv(XWX)
        except np.linalg.LinAlgError: XWX_inv = np.linalg.pinv(XWX)
        S = (np.sqrt(W)[:,None]) * X
        H = S @ XWX_inv @ S.T; h = np.clip(np.diag(H), 0, 1)
        U = X.T @ (y - p_hat + (0.5 - p_hat)*h)
        delta = XWX_inv @ U
        beta_new = beta + delta
        if np.linalg.norm(delta) < tol: beta = beta_new; break
        beta = beta_new
    eta = X @ beta; p_hat = 1/(1+np.exp(-eta)); W = np.clip(p_hat*(1-p_hat), 1e-9, None)
    XWX = X.T * W @ X
    try: cov = np.linalg.inv(XWX)
    except np.linalg.LinAlgError: cov = np.linalg.pinv(XWX)
    se = np.sqrt(np.clip(np.diag(cov), 1e-12, None))
    return beta, se, cov

def fit_logit(df, y_col, x_cols, robust=True, penalty=None, alpha=1.0, stepwise=False, firth=False):
    X = df[x_cols].copy()
    # ë²”ì£¼í˜• ìë™ ë”ë¯¸í™”(ë²”ì£¼ê°œìˆ˜<=6ë„ ë”ë¯¸)
    cat_cols=[c for c in x_cols if (not pd.api.types.is_numeric_dtype(X[c])) or X[c].nunique()<=6]
    X=pd.get_dummies(X, columns=cat_cols, drop_first=True, dtype=float)
    y=df[y_col].astype(int)
    data = pd.concat([y, X], axis=1).dropna()
    y = data[y_col]; X = data.drop(columns=[y_col])
    if X.shape[0] < X.shape[1] + 1 and not firth and penalty is None:
        raise ValueError("ìƒ˜í”Œ ìˆ˜ < ë³€ìˆ˜ ìˆ˜ (penalty ë˜ëŠ” Firth ì‚¬ìš© ê¶Œì¥)")
    if y.nunique()<2: raise ValueError("ì¢…ì†ë³€ìˆ˜ì— ë‘ ìˆ˜ì¤€ì´ í•„ìš”")
    X = sm.add_constant(X, has_constant="add")
    if firth:
        beta, se, cov = firth_logit(X.values, y.values)
        params = pd.Series(beta, index=X.columns)
        se_s = pd.Series(se, index=X.columns)
        z = params / se_s.replace(0, np.nan)
        pvals = 2*(1 - stats.norm.cdf(np.abs(z)))
        out = pd.DataFrame({"Variable":params.index, "OR":np.exp(params),
                            "CI Lower":np.exp(params - 1.96*se_s),
                            "CI Upper":np.exp(params + 1.96*se_s),
                            "P-value":pvals}).round(4)
        out = out[out["Variable"]!="const"].reset_index(drop=True)
        # VIF
        if HAS_VIF and X.shape[1]>1:
            Xv = X.drop(columns=["const"], errors="ignore")
            vif_rows = [{"Variable": Xv.columns[i], "VIF": variance_inflation_factor(Xv.values, i)} for i in range(Xv.shape[1])] if Xv.shape[1]>=2 else []
        else:
            vif_rows = []
        class Res: pass
        res = Res(); res.params=params; res.cov_params=lambda: cov; res.pvalues=pd.Series(pvals, index=params.index)
        res.summary2=lambda: type("S",(),{"as_text":lambda self: "Firth logistic (approx Wald CIs)"})()
        return res, out, pd.DataFrame(vif_rows)
    if stepwise and penalty is None:
        current=[]; remaining=list(X.columns); remaining.remove("const"); best_aic=np.inf
        while True:
            scores=[]; improved=False
            for v in remaining:
                trial=["const"]+current+[v]
                try: m=sm.Logit(y, X[trial]).fit(disp=False, maxiter=200); aic=m.aic; scores.append((aic,v,m))
                except Exception: continue
            if not scores: break
            scores.sort(key=lambda x:x[0])
            if scores[0][0]+1e-6 < best_aic:
                best_aic=scores[0][0]; current=current+[scores[0][1]]; remaining.remove(scores[0][1]); improved=True
            if not improved: break
        model = sm.Logit(y, X[["const"]+current]).fit(disp=False, maxiter=200) if current else sm.Logit(y,X).fit(disp=False, maxiter=200)
        res = model.get_robustcov_results(cov_type="HC3") if robust else model
        conf = res.conf_int()
        out = pd.DataFrame({"Variable":res.params.index, "OR":np.exp(res.params),
                            "CI Lower":np.exp(conf[0]), "CI Upper":np.exp(conf[1]), "P-value":res.pvalues}).round(4)
        out = out[out["Variable"]!="const"].reset_index(drop=True)
        # VIF
        if HAS_VIF:
            Xv = (X[["const"]+current]).drop(columns=["const"], errors="ignore")
            vif_rows = [{"Variable": Xv.columns[i], "VIF": variance_inflation_factor(Xv.values, i)} for i in range(Xv.shape[1])] if Xv.shape[1]>=2 else []
        else:
            vif_rows = []
        return res, out, pd.DataFrame(vif_rows)
    if penalty in ("l2","l1"):
        l1_wt=0.0 if penalty=="l2" else 1.0
        m=sm.Logit(y, X).fit_regularized(alpha=alpha, L1_wt=l1_wt, disp=False)
        params=m.params
        out=pd.DataFrame({"Variable":params.index, "OR":np.exp(params)}).round(4)
        out=out[out["Variable"]!="const"].reset_index(drop=True)
        if HAS_VIF:
            Xv=X.drop(columns=["const"], errors="ignore")
            vif_rows=[{"Variable":Xv.columns[i], "VIF":variance_inflation_factor(Xv.values,i)} for i in range(Xv.shape[1])] if Xv.shape[1]>=2 else []
        else:
            vif_rows=[]
        return m, out, pd.DataFrame(vif_rows)
    model = sm.Logit(y, X).fit(disp=False, maxiter=200)
    res = model.get_robustcov_results(cov_type="HC3") if robust else model
    conf = res.conf_int()
    out = pd.DataFrame({"Variable":res.params.index, "OR":np.exp(res.params),
                        "CI Lower":np.exp(conf[0]), "CI Upper":np.exp(conf[1]), "P-value":res.pvalues}).round(4)
    out = out[out["Variable"]!="const"].reset_index(drop=True)
    if HAS_VIF:
        Xv = X.drop(columns=["const"], errors="ignore")
        vif_rows = [{"Variable": Xv.columns[i], "VIF": variance_inflation_factor(Xv.values, i)} for i in range(Xv.shape[1])] if Xv.shape[1]>=2 else []
    else:
        vif_rows=[]
    return res, out, pd.DataFrame(vif_rows)

def infer_labels(raters_df):
    vals = [v for v in pd.unique(raters_df.values.ravel()) if pd.notna(v)]
    return sorted(vals) if vals else [1,2,3,4,5]

def per_rater_stats(raters_df, labels, D):
    raters = list(raters_df.columns)
    means=[]
    for r in raters:
        ks=[]
        for s in raters:
            if r==s: continue
            ks.append(weighted_kappa_custom(raters_df[r], raters_df[s], labels, D))
        means.append(np.nanmean(ks))
    miss = raters_df.isna().mean().values*100.0
    return pd.DataFrame({"Rater":raters,"Mean weighted Îº":np.round(means,4),"Missing (%)":np.round(miss,2)})

def kappa_heatmap(raters_df, labels, D, title):
    raters = list(raters_df.columns)
    K = np.ones((len(raters), len(raters)), float)
    for i,a in enumerate(raters):
        for j,b in enumerate(raters):
            if i<j:
                k = weighted_kappa_custom(raters_df[a], raters_df[b], labels, D)
                K[i,j]=K[j,i]=k
    if not HAS_PLOTLY:
        return None, K, raters
    fig = px.imshow(K, x=raters, y=raters, zmin=-1, zmax=1, text_auto=".2f", title=title, aspect="auto")
    fig.update_layout(coloraxis_colorbar=dict(title="Îº"), xaxis_title="", yaxis_title="")
    return fig, K, raters

def plot_roc_curves(y_true, scores_dict, auc_dict):
    fig = go.Figure()
    fig.add_shape(type='line', line=dict(dash='dash'), x0=0, x1=1, y0=0, y1=1)
    for name, score in scores_dict.items():
        idx = y_true.dropna().index.intersection(score.dropna().index)
        yt = y_true.loc[idx]; ys = score.loc[idx]
        if len(idx)==0: continue
        fpr, tpr, _ = roc_curve(yt, ys)
        auc_val = auc_dict.get(name, np.nan)
        fig.add_trace(go.Scatter(x=fpr, y=tpr, name=f"{name} (AUC={auc_val:.3f})", mode='lines'))
    fig.update_layout(
        xaxis_title='1 - Specificity',
        yaxis_title='Sensitivity',
        title='ROC Curve Comparison',
        legend=dict(x=0.5, y=0.1, xanchor='center', yanchor='bottom')
    )
    return fig

def summarize_pair(name, raters_df, weight_scheme, custom_D, B, seed, kripp_level="ordinal"):
    labels = infer_labels(raters_df)
    D = custom_D if custom_D is not None else make_distance_matrix(len(labels), weight_scheme)
    pair_tab = pairwise_kappa_table(raters_df, labels, D, B=B, seed=seed)
    lk,llo,lhi,_ = lights_kappa(raters_df, labels, D, B=B, seed=seed)
    fk, fk_lo, fk_hi, _ = fleiss_kappa(raters_df, B=B, seed=seed) if raters_df.shape[1]>=2 else (np.nan,np.nan,np.nan,np.array([]))
    ka, ka_lo, ka_hi, _ = krippendorff_alpha(raters_df, level=kripp_level, B=B, seed=seed)
    icc, icc_lo, icc_hi = calculate_icc(raters_df)
    # ë³´ìˆ˜ì : BP/PABAKì—ë„ ë¶€íŠ¸ìŠ¤íŠ¸ë© CI
    bp = brennan_prediger_kappa(raters_df)
    bp_lo, bp_hi, _ = bootstrap_multi_stat(raters_df, brennan_prediger_kappa, B=B, seed=seed)
    pb = pabak_multi(raters_df)
    pb_lo, pb_hi, _ = bootstrap_multi_stat(raters_df, pabak_multi, B=B, seed=seed)

    summary = pd.DataFrame([
        {"Metric":"Light's Îº (weighted mean)","Value":lk,"CI_lo":llo,"CI_hi":lhi,"Interp":interpret_kappa(lk),"Definition":GLOSSARY["Light's Îº (weighted mean)"]},
        {"Metric":"Fleiss' Îº (unweighted)","Value":fk,"CI_lo":fk_lo,"CI_hi":fk_hi,"Interp":interpret_kappa(fk),"Definition":GLOSSARY["Fleiss' Îº (unweighted)"]},
        {"Metric":f"Krippendorff's Î± ({kripp_level})","Value":ka,"CI_lo":ka_lo,"CI_hi":ka_hi,"Interp":interpret_kappa(ka),"Definition":GLOSSARY["Krippendorff's Î± (ordinal)"]},
        {"Metric":"ICC(2,k)","Value":icc,"CI_lo":icc_lo,"CI_hi":icc_hi,"Interp":interpret_kappa(icc),"Definition":GLOSSARY["ICC(2,k)"]},
        {"Metric":"Brennanâ€“Prediger Îº","Value":bp,"CI_lo":bp_lo,"CI_hi":bp_hi,"Interp":interpret_kappa(bp),"Definition":GLOSSARY["Brennanâ€“Prediger Îº"]},
        {"Metric":"PABAK (multi)","Value":pb,"CI_lo":pb_lo,"CI_hi":pb_hi,"Interp":interpret_kappa(pb),"Definition":GLOSSARY["PABAK (multi)"]},
    ]).round(4)

    consensus = raters_df.median(axis=1, skipna=True)
    rat_stats = per_rater_stats(raters_df, labels, D)
    return pair_tab, summary, consensus, labels, D, rat_stats

# ---------------- Sidebar ----------------
with st.sidebar:
    st.header("1) íŒŒì¼ ì—…ë¡œë“œ")
    up = st.file_uploader("Excel (.xlsx)", type=["xlsx"], key="uploader_main")
    if st.button("ìƒ˜í”Œ í…œí”Œë¦¿(.xlsx) ë‹¤ìš´ë¡œë“œ", key="btn_sample"):
        bio=BytesIO(); st.session_state["_dummy"]=0
        np.random.seed(0); n=300
        def r5(n): return np.random.choice([1,2,3,4,5], size=n, p=[.1,.2,.4,.2,.1])
        raters=6
        cols={}
        for pair in ["A_M","A_G","M_G"]:
            for r in range(1, raters+1): cols[f"{pair}_S{r}"]=r5(n)
        df_tmp=pd.DataFrame({"case_id":[f"C{str(i+1).zfill(3)}" for i in range(n)], **cols,
                             "AI_label": r5(n), "MDT_label": r5(n), "GL_label": r5(n),
                             "age":np.random.normal(65,10,n).round().astype(int),
                             "sex":np.random.choice(["M","F"],n),
                             "ASA":np.random.choice([1,2,3,4],n,p=[.1,.4,.4,.1]),
                             "ECOG":np.random.choice([0,1,2,3],n,p=[.3,.4,.2,.1]),
                             "T":np.random.choice(["T1","T2","T3","T4"],n,p=[.1,.2,.5,.2]),
                             "N":np.random.choice(["N0","N1","N2"],n,p=[.5,.3,.2]),
                             "stage":np.random.choice(["I","II","III","IV"],n,p=[.1,.4,.4,.1]),
                             "PNI":np.random.choice([0,1],n,p=[.7,.3]),
                             "EMVI":np.random.choice([0,1],n,p=[.7,.3]),
                             "obstruction":np.random.choice([0,1],n,p=[.7,.3])})
        with pd.ExcelWriter(bio, engine="openpyxl") as w: df_tmp.to_excel(w, index=False)
        bio.seek(0)
        st.download_button("sample_template.xlsx ì €ì¥", data=bio, file_name="sample_template.xlsx", key="dl_sample")

    st.divider()
    st.header("2) ë¶„ì„ ì„¤ì •")
    B = st.number_input("Bootstrap resamples", value=2000, min_value=200, step=200, key="num_B", help="ë¶€íŠ¸ìŠ¤íŠ¸ë© ë°˜ë³µ ìˆ˜(í´ìˆ˜ë¡ ì‹ ë¢°êµ¬ê°„ ì•ˆì •).")
    seed = st.number_input("Random seed", value=42, min_value=0, step=1, key="num_seed", help="ì¬í˜„ì„±ì„ ìœ„í•œ ë‚œìˆ˜ ì‹œë“œ.")
    weight_scheme = st.selectbox("ê°€ì¤‘ ë°©ì‹ (Îº)", ["quadratic","linear","stepwise","unweighted","custom CSV"], key="sel_weight",
                                 help="ê°€ì¤‘ Cohen Îºì—ì„œ ë²”ì£¼ ê°„ ê±°ë¦¬ í•¨ìˆ˜ ì„ íƒ.")
    custom_D = None
    if weight_scheme == "custom CSV":
        up_csv = st.file_uploader("ê°€ì¤‘ ê±°ë¦¬í–‰ë ¬ CSV ì—…ë¡œë“œ (kÃ—k, ëŒ€ê°=0, [0,1])", type=["csv"], key="W_csv")
    else:
        up_csv = None
    kripp_level = st.selectbox("Krippendorff's Î± ë“±ê¸‰", ["ordinal","nominal"], index=0, key="sel_alpha_level",
                               help="Î± ê±°ë¦¬: ordinal=ì„œì—´(ì œê³±ê±°ë¦¬), nominal=ëª…ëª©(0/1).")
    st.divider()
    st.header("3) í¬í•¨í•  ë¶„ì„/ì˜µì…˜")
    do_icc    = st.checkbox("ICC í¬í•¨", value=True, key="chk_icc", help=GLOSSARY["ICC(2,k)"])
    do_fleiss = st.checkbox("Fleiss' Îº í¬í•¨", value=True, key="chk_fleiss", help=GLOSSARY["Fleiss' Îº (unweighted)"])
    do_kripp  = st.checkbox("Krippendorff's Î± í¬í•¨", value=True, key="chk_kripp", help=GLOSSARY["Krippendorff's Î± (ordinal)"])
    do_mcnemar= st.checkbox("McNemar (ìŒ ë¹„êµ)", value=True, key="chk_mcnemar", help=GLOSSARY["McNemar"])
    do_auc    = st.checkbox("AUC/ì •í™•ë„ ë¹„êµ + DeLong", value=True, key="chk_auc", help=f"{GLOSSARY['AUC']} {GLOSSARY['DeLong']}")
    do_logit  = st.checkbox("ë¶ˆì¼ì¹˜ ìœ„í—˜ìš”ì¸(ë¡œì§€ìŠ¤í‹±)", value=True, key="chk_logit", help=GLOSSARY["Firth logistic"])
    do_direct_kappa = st.checkbox("ì§ì ‘ Cohen weighted Îº (AI_label/MDT_label/GL_label)", value=True, key="chk_direct")
    export_all_zip  = st.checkbox("ëª¨ë“  ì‚°ì¶œë¬¼ ZIP ë¬¶ìŒ ì œê³µ", value=True, key="chk_zip")
    st.divider()
    run_all = st.button("ğŸš€ í•œ ë²ˆì— ì „ì²´ ì‹¤í–‰", type="primary", use_container_width=True, key="btn_runall")

# ---------------- Load data ----------------
if up is None:
    st.info("ì—‘ì…€ì„ ì—…ë¡œë“œí•˜ì„¸ìš”. ë˜ëŠ” ìƒ˜í”Œ í…œí”Œë¦¿ì„ ì‚¬ìš©í•˜ì„¸ìš”.")
    st.stop()

df = read_excel(up)
if df is None:
    st.stop()
st.success(f"Loaded: {df.shape[0]} rows Ã— {df.shape[1]} cols")

with st.expander("ë°ì´í„° ë¯¸ë¦¬ë³´ê¸°", expanded=False):
    st.dataframe(df.head(), use_container_width=True)

# ---------------- Column mapping ----------------
with st.expander("ì—´ ì„ íƒ â€” í‰ê°€ì ìˆ˜ ììœ ", expanded=True):
    cols = df.columns.tolist()
    case_id_col = st.selectbox("Case ID ì—´", options=cols, index=cols.index("case_id") if "case_id" in cols else 0, key="sel_caseid")
    st.markdown("**AI vs MDT**")
    am_cols = st.multiselect("í‰ê°€ì ì—´(2ê°œ ì´ìƒ) â€” AI vs MDT", options=cols, default=[c for c in cols if c.startswith("A_M_")], key="ms_am")
    st.markdown("**AI vs Guideline**")
    ag_cols = st.multiselect("í‰ê°€ì ì—´(2ê°œ ì´ìƒ) â€” AI vs Guideline", options=cols, default=[c for c in cols if c.startswith("A_G_")], key="ms_ag")
    st.markdown("**MDT vs Guideline**")
    mg_cols = st.multiselect("í‰ê°€ì ì—´(2ê°œ ì´ìƒ) â€” MDT vs Guideline", options=cols, default=[c for c in cols if c.startswith("M_G_")], key="ms_mg")

for c in am_cols+ag_cols+mg_cols:
    if c in df.columns: df[c]=safe_num(df[c])

W_custom = None
if weight_scheme=="custom CSV" and up_csv is not None:
    try:
        W_custom = sanitize_custom_D(pd.read_csv(up_csv, header=None).values)
    except Exception as e:
        st.error(f"ê°€ì¤‘í–‰ë ¬ CSV ì˜¤ë¥˜: {e}")
        W_custom = None

# ---------------- Run all ----------------
if run_all:
    with st.spinner("ë¶„ì„ì„ ì‹¤í–‰ ì¤‘ì…ë‹ˆë‹¤..."):
        # ë§¤ ì‹¤í–‰ë§ˆë‹¤ ê²°ê³¼ ì´ˆê¸°í™”(ì„¸ì…˜ ìƒíƒœì— ì €ì¥)
        st.session_state.results_tables = {}
        st.session_state.consensus_map = {}
        st.session_state.images_store = {}

        pairs = {"AI_vs_MDT": am_cols, "AI_vs_Guideline": ag_cols, "MDT_vs_Guideline": mg_cols}
        for name, cols_list in pairs.items():
            st.header(f"{name}")
            if len(cols_list)<2:
                st.warning(f"{name}: í‰ê°€ì ì—´ì„ 2ê°œ ì´ìƒ ì„ íƒí•˜ì„¸ìš”.")
                continue
            R = df[cols_list].copy()
            labels = sorted([v for v in pd.unique(R.values.ravel()) if pd.notna(v)]) or [1,2,3,4,5]
            if W_custom is not None and W_custom.shape != (len(labels),len(labels)):
                st.error(f"{name}: ì»¤ìŠ¤í…€ ê°€ì¤‘í–‰ë ¬ í¬ê¸° {W_custom.shape} ì´ ë¼ë²¨ ìˆ˜ k={len(labels)}ì™€ ì¼ì¹˜í•˜ì§€ ì•ŠìŠµë‹ˆë‹¤. ê¸°ë³¸ {weight_scheme if weight_scheme!='custom CSV' else 'quadratic'} ì‚¬ìš©")
                D_use = None
            else:
                D_use = W_custom
            pair_tab, summary, consensus, labels_used, D_used, rat_stats = summarize_pair(name, R, weight_scheme, D_use, int(B), int(seed), kripp_level=kripp_level)

            if not do_icc:
                summary = summary[summary["Metric"]!="ICC(2,k)"]
            if not do_fleiss:
                summary = summary[summary["Metric"]!="Fleiss' Îº (unweighted)"]
            if not do_kripp:
                summary = summary[~summary["Metric"].str.startswith("Krippendorff")]

            tab1, tab2, tab3, tab4 = st.tabs(["ìŒë³„ ê°€ì¤‘ Îº", "ìš”ì•½ ì§€í‘œ", "í‰ê°€ìë³„ í†µê³„", "ì‹œê°í™”"])
            with tab1:
                st.subheader("ìŒë³„ ê°€ì¤‘ Îº (+Holm ë³´ì •)")
                st.dataframe(pair_tab, use_container_width=True)
            with tab2:
                st.subheader("ìš”ì•½ ì§€í‘œ")
                st.dataframe(summary, use_container_width=True)
            with tab3:
                st.subheader("í‰ê°€ìë³„ í†µê³„ (í‰ê·  Îº, ê²°ì¸¡ë¥ )")
                st.dataframe(rat_stats, use_container_width=True)

            st.session_state.results_tables[f"{name}_pairwise_kappa"]=pair_tab
            st.session_state.results_tables[f"{name}_summary"]=summary
            st.session_state.results_tables[f"{name}_per_rater"]=rat_stats
            st.session_state.consensus_map[name]=consensus

            if HAS_PLOTLY:
                with tab4:
                    st.subheader("í•©ì˜ ì ìˆ˜ ë¶„í¬")
                    fig_hist=px.histogram(consensus.dropna(), nbins=max(10, len(labels_used)*2), title=f"{name} consensus histogram")
                    fig_hist.update_layout(showlegend=False); st.plotly_chart(fig_hist, use_container_width=True)
                    try:
                        bio=BytesIO(); fig_hist.write_image(bio, format="png"); bio.seek(0); st.session_state.images_store[f"{name}_consensus_hist.png"]=bio.getvalue()
                    except Exception: pass

                    st.subheader("ìŒë³„ Îº íˆíŠ¸ë§µ")
                    fig_h, K, raters = kappa_heatmap(R, labels_used, D_used if D_use is not None else make_distance_matrix(len(labels_used), weight_scheme), f"{name} pairwise weighted Îº heatmap")
                    if fig_h is not None:
                        st.plotly_chart(fig_h, use_container_width=True)
                        try:
                            bio=BytesIO(); fig_h.write_image(bio, format="png"); bio.seek(0); st.session_state.images_store[f"{name}_kappa_heatmap.png"]=bio.getvalue()
                        except Exception: pass

    # Direct Cohen weighted Îº for raw labels
    if do_direct_kappa and all(c in df.columns for c in ["AI_label","MDT_label","GL_label"]):
        st.header("ì§ì ‘ Cohen weighted Îº (ì›ì‹œ ë ˆì´ë¸”)")
        labels = sorted([v for v in pd.unique(df[["AI_label","MDT_label","GL_label"]].values.ravel()) if pd.notna(v)]) or [1,2,3,4,5]
        D = W_custom if (W_custom is not None and W_custom.shape==(len(labels),len(labels))) else make_distance_matrix(len(labels), weight_scheme if weight_scheme!="custom CSV" else "quadratic")
        co_tab = []
        for (a,b,name2) in [("AI_label","MDT_label","AI_vs_MDT_raw"),
                             ("AI_label","GL_label","AI_vs_GL_raw"),
                             ("MDT_label","GL_label","MDT_vs_GL_raw")]:
            k0 = weighted_kappa_custom(df[a], df[b], labels, D)
            def stat(idx): return weighted_kappa_custom(df[a].iloc[idx], df[b].iloc[idx], labels, D)
            lo,hi,p,_ = bootstrap_ci_stat(stat, len(df), B=int(B), seed=int(seed))
            co_tab.append({"Pair":name2,"Weighted Kappa":round(k0,4),"95% CI Lower":round(lo,4),"95% CI Upper":round(hi,4),"P (bootstrap)":round(p,6),"Interp":interpret_kappa(k0)})
        co_df = pd.DataFrame(co_tab)
        st.dataframe(co_df, use_container_width=True)
        st.session_state.results_tables["Direct_Cohen_weighted_kappa"]=co_df

    # ------------- McNemar & AUC -------------
    if do_mcnemar or do_auc:
        st.header("ì„±ëŠ¥ ë¹„êµ (McNemar, AUC/ì •í™•ë„, DeLong)")
        all_keys = list(st.session_state.consensus_map.keys())
        if len(all_keys)>=2:
            left_key = st.selectbox("ì˜ˆì¸¡1", all_keys, index=0, key="pred1_key")
            right_key = st.selectbox("ì˜ˆì¸¡2", all_keys, index=1, key="pred2_key")
            thr1 = st.number_input(f"{left_key} ì–‘ì„± ì„ê³„(â‰¥)", value=4, min_value=1, max_value=5, step=1, key="thr1")
            thr2 = st.number_input(f"{right_key} ì–‘ì„± ì„ê³„(â‰¥)", value=4, min_value=1, max_value=5, step=1, key="thr2")
            bin1 = (st.session_state.consensus_map[left_key] >= thr1).astype(float)
            bin2 = (st.session_state.consensus_map[right_key] >= thr2).astype(float)
            v = ~(bin1.isna() | bin2.isna())
            tab = pd.crosstab(bin1[v], bin2[v])
            st.dataframe(tab, use_container_width=True)
            st.session_state.results_tables["McNemar_table"]=tab.reset_index()
            if do_mcnemar and tab.shape==(2,2):
                try:
                    res = sm_mcnemar(tab.to_numpy(), exact=True)
                    st.write(f"**McNemar exact p**: {res.pvalue:.6f}")
                    st.session_state.results_tables["McNemar_p"]=pd.DataFrame({"p_exact":[res.pvalue]})
                except Exception as e:
                    st.error(f"McNemar ì˜¤ë¥˜: {e}")

            if do_auc:
                ref_key = st.selectbox("Reference(ì´ì§„)ë¡œ ì‚¬ìš©í•  ìŒ", all_keys, index=0, key="ref_key_auc")
                thr_ref = st.number_input(f"{ref_key} ê¸°ì¤€ ì„ê³„(â‰¥)", value=4, min_value=1, max_value=5, step=1, key="thr_ref")
                y_true = (st.session_state.consensus_map[ref_key] >= thr_ref).astype(int)
                pred_keys = st.multiselect("ì˜ˆì¸¡ìë¡œ ë¹„êµí•  ìŒ(ì ìˆ˜)", all_keys, default=[k for k in all_keys if k!=ref_key], key="pred_keys_auc")
                rows=[]; ref_for_delong=None; scores_for_roc = {}; aucs_for_roc = {}
                for pk in pred_keys:
                    score = (st.session_state.consensus_map[pk] - st.session_state.consensus_map[pk].min()) / max(1e-9, (st.session_state.consensus_map[pk].max()-st.session_state.consensus_map[pk].min()))
                    idx = y_true.dropna().index.intersection(score.dropna().index)
                    if len(idx)==0:
                        rows.append({"Predictor":pk,"AUC":np.nan,"AUC_CI_lo":np.nan,"AUC_CI_hi":np.nan,"DeLong_p":np.nan,
                                     "Accuracy":np.nan,"Acc_CI_lo":np.nan,"Acc_CI_hi":np.nan})
                        continue
                    yt = y_true.loc[idx]; ys = score.loc[idx]
                    auc = roc_auc_score(yt, ys)
                    lo, hi, _ = bootstrap_metric_ci(yt, ys, lambda u,v: roc_auc_score(u,v), B=int(B), seed=int(seed))
                    scores_for_roc[pk] = ys; aucs_for_roc[pk] = auc
                    if ref_for_delong is None:
                        ref_for_delong = ys; delong_p = np.nan
                    else:
                        try:
                            _,_,_,delong_p = delong_compare(yt.values, ref_for_delong.values, ys.values)
                        except Exception:
                            delong_p = np.nan
                    cutoff = st.number_input(f"{pk} ì´ì§„í™” ì„ê³„(ì •í™•ë„)", value=4, min_value=1, max_value=5, step=1, key=f"cut_{pk}")
                    acc = float(( (st.session_state.consensus_map[pk].loc[idx] >= cutoff).astype(float) == yt ).mean())
                    lo_a, hi_a, _ = bootstrap_metric_ci(yt, (st.session_state.consensus_map[pk].loc[idx] >= cutoff).astype(float),
                                                        lambda u,v: float((u==v).mean()), B=int(B), seed=int(seed))
                    rows.append({"Predictor":pk,"AUC":round(auc,4),"AUC_CI_lo":round(lo,4),"AUC_CI_hi":round(hi,4),
                                 "DeLong_p": None if np.isnan(delong_p) else round(delong_p,6),
                                 "Accuracy":round(acc,4),"Acc_CI_lo":round(lo_a,4),"Acc_CI_hi":round(hi_a,4)})
                perf_df = pd.DataFrame(rows)
                st.dataframe(perf_df, use_container_width=True)
                st.session_state.results_tables["Performance"]=perf_df
                if scores_for_roc and HAS_PLOTLY:
                    fig_roc = plot_roc_curves(y_true, scores_for_roc, aucs_for_roc)
                    st.plotly_chart(fig_roc, use_container_width=True)
                    try:
                        bio = BytesIO(); fig_roc.write_image(bio, format="png"); bio.seek(0)
                        st.session_state.images_store["ROC_comparison.png"] = bio.getvalue()
                    except Exception: pass

    # ------------- Logistic risk factors -------------
    if do_logit:
        st.header("ë¶ˆì¼ì¹˜ ìœ„í—˜ìš”ì¸ (ë¡œì§€ìŠ¤í‹±)")
        model_pairs = st.multiselect("íšŒê·€ë¶„ì„ì„ ìˆ˜í–‰í•  ë¹„êµìŒ ì„ íƒ", list(st.session_state.consensus_map.keys()), default=list(st.session_state.consensus_map.keys()), key="ms_models")
        thr_discord = st.number_input("ë¶ˆì¼ì¹˜ ì •ì˜: í•©ì˜ ì ìˆ˜ â‰¤", value=2, min_value=1, max_value=5, step=1, key="num_thr_discord", help="í•©ì˜ ì¤‘ì•™ê°’ì´ ì„ê³„ ì´í•˜ì´ë©´ ë¶ˆì¼ì¹˜ë¡œ ì •ì˜.")
        excluded = set([case_id_col]) | set(am_cols) | set(ag_cols) | set(mg_cols)
        candidates = [c for c in df.columns if c not in excluded]
        default_covars = [c for c in ["age","sex","ASA","ECOG","stage","T","N","PNI","EMVI","obstruction"] if c in candidates]
        covars = st.multiselect("ê³µë³€ëŸ‰ ì„ íƒ", options=candidates, default=default_covars, key="ms_covars")
        robust = st.checkbox("Robust SE (HC3)", value=True, key="chk_robust", help="ì´ë¶„ì‚°Â·ëª¨í˜•ì˜¤ë¥˜ì— ê²¬ê³ í•œ í‘œì¤€ì˜¤ì°¨.")
        penalty = st.selectbox("íŒ¨ë„í‹°(ì„ íƒ)", ["none","l2","l1"], index=0, key="sel_penalty", help="ì •ê·œí™” íšŒê·€(L2/L1). ìƒ˜í”Œ ëŒ€ë¹„ ë³€ìˆ˜ ë§ì„ ë•Œ ê¶Œì¥.")
        alpha = st.number_input("íŒ¨ë„í‹° ê°•ë„ alpha", value=1.0, min_value=0.0, step=0.1, key="num_alpha")
        stepwise = st.checkbox("Stepwise (AIC) ì‚¬ìš© (íŒ¨ë„í‹° ë¯¸ì‚¬ìš© ì‹œ)", value=False, key="chk_stepwise")
        use_firth = st.checkbox("Firth (ë¶„ë¦¬/í¬ê·€ì‚¬ê±´ ëŒ€ë¹„)", value=False, key="chk_firth", help=GLOSSARY["Firth logistic"])
        for pair in model_pairs:
            if pair not in st.session_state.consensus_map:
                continue
            y = (st.session_state.consensus_map[pair] <= thr_discord).astype(int)
            df["_discordant"] = y
            try:
                pen = None if penalty=="none" else penalty
                res, or_tab, vif_df = fit_logit(df, "_discordant", covars, robust=robust, penalty=pen, alpha=alpha, stepwise=stepwise, firth=use_firth)
                st.subheader(f"{pair} â€” OR (discordance=1)")
                logit_tab1, logit_tab2 = st.tabs(["Odds Ratios", "VIF"])
                with logit_tab1:
                    st.dataframe(or_tab, use_container_width=True)
                with logit_tab2:
                    if HAS_VIF and not vif_df.empty:
                        st.dataframe(vif_df.sort_values("VIF", ascending=False), use_container_width=True)
                    elif not HAS_VIF:
                        st.info("statsmodelsì˜ VIF ëª¨ë“ˆì„ ì°¾ì§€ ëª»í–ˆìŠµë‹ˆë‹¤. (ì„ íƒ ì‚¬í•­) `pip install -U statsmodels` í›„ ì‚¬ìš© ê°€ëŠ¥í•©ë‹ˆë‹¤.")
                    else:
                        st.info("VIFë¥¼ ê³„ì‚°í•  ë³€ìˆ˜ê°€ 2ê°œ ë¯¸ë§Œì…ë‹ˆë‹¤.")
                st.session_state.results_tables[f"RiskFactors_{pair}"]=or_tab
                if not vif_df.empty:
                    st.session_state.results_tables[f"VIF_{pair}"]=vif_df
                try:
                    st.session_state.results_tables[f"ModelSummary_{pair}"]=pd.DataFrame({"summary":[res.summary2().as_text()]})
                except Exception:
                    st.session_state.results_tables[f"ModelSummary_{pair}"]=pd.DataFrame({"summary":[str(res)]})
            except Exception as e:
                st.error(f"{pair} íšŒê·€ ì˜¤ë¥˜: {e}")

# ---------------- Export: Excel + Reports + ALL ZIP ----------------
st.divider()
st.subheader("ë³´ê³ ì„œ/ê²°ê³¼ ë‹¤ìš´ë¡œë“œ")

def to_excel(tables: dict):
    bio = BytesIO()
    with pd.ExcelWriter(bio, engine="openpyxl") as w:
        for name, d in tables.items():
            if isinstance(d, pd.DataFrame) and not d.empty:
                d.to_excel(w, sheet_name=name[:31], index=False)
            elif isinstance(d, pd.DataFrame) and d.empty:
                continue
            else:
                if isinstance(d, str):
                    pd.DataFrame({"text":[d]}).to_excel(w, sheet_name=name[:31], index=False)
    bio.seek(0); return bio

def build_pptx(tables: dict):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        st.error(f"PPTX ë¼ì´ë¸ŒëŸ¬ë¦¬ ì˜¤ë¥˜: `pip install python-pptx` í•„ìš”. ì˜¤ë¥˜: {e}"); return None
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "CRC Agreement â€” Summary"; slide.placeholders[1].text = "Auto-generated report"
    def add_table_slide(title, df):
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = title[:50]
        rows, cols = df.shape
        table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
        for j,c in enumerate(df.columns): table.cell(0,j).text = str(c)
        for i in range(rows):
            for j in range(cols):
                table.cell(i+1,j).text = str(df.iloc[i,j])
    for name, d in tables.items():
        if isinstance(d, pd.DataFrame) and not d.empty:
            add_table_slide(name, d.iloc[:30, :10])
    bio=BytesIO(); prs.save(bio); bio.seek(0); return bio

def build_docx(tables: dict):
    try:
        from docx import Document
    except Exception as e:
        st.error(f"DOCX ë¼ì´ë¸ŒëŸ¬ë¦¬ ì˜¤ë¥˜: `pip install python-docx` í•„ìš”. ì˜¤ë¥˜: {e}"); return None
    doc = Document(); doc.add_heading("CRC Agreement â€” Summary", 0)
    for name, d in tables.items():
        doc.add_heading(name, level=1)
        if isinstance(d, pd.DataFrame) and not d.empty:
            t = doc.add_table(rows=d.shape[0]+1, cols=d.shape[1])
            t.style = 'Table Grid'
            for j,c in enumerate(d.columns): t.cell(0,j).text = str(c)
            for i in range(d.shape[0]):
                for j in range(d.shape[1]):
                    # âœ¨ ì—¬ê¸°ê°€ ìˆ˜ì •ëœ ë¶€ë¶„ì…ë‹ˆë‹¤ âœ¨
                    t.cell(i+1,j).text = str(d.iloc[i,j])
        else:
            doc.add_paragraph(str(d))
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio

def build_tex(tables: dict):
    def esc(s):
        return (str(s)
                .replace('\\', r'\textbackslash{}')
                .replace('_', r'\_')
                .replace('%', r'\%')
                .replace('&', r'\&')
                .replace('#', r'\#'))
    lines = [
        r"\documentclass{article}",
        r"\usepackage{booktabs}",
        r"\usepackage[margin=1in]{geometry}",
        r"\begin{document}",
        r"\section*{CRC Agreement --- Summary}",
    ]
    for name, d in tables.items():
        lines.append(r"\subsection*{" + esc(name) + r"}")
        if isinstance(d, pd.DataFrame) and not d.empty:
            lines.append(r"\begin{tabular}{" + ("l"*d.shape[1]) + r"}")
            header = " & ".join(map(esc, d.columns)) + r" \\"
            lines.append(r"\toprule")
            lines.append(header)
            lines.append(r"\midrule")
            for i in range(d.shape[0]):
                row = " & ".join([esc(x) for x in d.iloc[i,:].tolist()]) + r" \\"
                lines.append(row)
            lines.append(r"\bottomrule")
            lines.append(r"\end{tabular}")
        else:
            lines.append(esc(d))
    lines.append(r"\end{document}")
    bio=BytesIO(); bio.write(("\n".join(lines)).encode("utf-8")); bio.seek(0); return bio

def pack_all_zip(tables: dict, images: dict):
    zbio = BytesIO()
    with ZipFile(zbio, "w", ZIP_DEFLATED) as z:
        xls = to_excel(tables); z.writestr("results.xlsx", xls.getvalue())
        pptx = build_pptx(tables)
        if pptx is not None: z.writestr("report.pptx", pptx.getvalue())
        docx = build_docx(tables)
        if docx is not None: z.writestr("report.docx", docx.getvalue())
        tex = build_tex(tables); z.writestr("report.tex", tex.getvalue())
        for name, data in images.items():
            z.writestr(f"figures/{name}", data)
        try:
            j = {k: (v.to_dict(orient="list") if isinstance(v, pd.DataFrame) else str(v)) for k,v in tables.items()}
            z.writestr("results.json", json.dumps(j, indent=2).encode("utf-8"))
        except Exception:
            pass
    zbio.seek(0); return zbio

# --- ì•ˆì „í•˜ê²Œ ì„¸ì…˜ì—ì„œ ì½ê¸° ---
results_tables = st.session_state.get("results_tables", {})
images_store  = st.session_state.get("images_store", {})

if results_tables:
    xls = to_excel(results_tables)
    st.download_button("í†µí•© ê²°ê³¼ Excel (results.xlsx)", data=xls, file_name="results.xlsx",
                      mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", key="dl_xlsx")
    pptx = build_pptx(results_tables)
    if pptx is not None:
        st.download_button("PPTX ë³´ê³ ì„œ", data=pptx, file_name="report.pptx",
                          mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", key="dl_pptx")
    docx = build_docx(results_tables)
    if docx is not None:
        st.download_button("DOCX ë³´ê³ ì„œ", data=docx, file_name="report.docx",
                          mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", key="dl_docx")
    tex = build_tex(results_tables)
    st.download_button("LaTeX (.tex)", data=tex, file_name="report.tex", mime="application/x-tex", key="dl_tex")
    if export_all_zip:
        allzip = pack_all_zip(results_tables, images_store)
        st.download_button("ğŸ“¦ ëª¨ë“  ì‚°ì¶œë¬¼ ZIP (results_all.zip)", data=allzip, file_name="results_all.zip", mime="application/zip", key="dl_zip")
else:
    st.info("ì•„ì§ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤. ì‚¬ì´ë“œë°”ì˜ 'í•œ ë²ˆì— ì „ì²´ ì‹¤í–‰'ì„ ëˆŒëŸ¬ì£¼ì„¸ìš”.")

# ---------------- Glossary section ----------------
with st.expander("â„¹ï¸ ìš©ì–´ ì„¤ëª…(ê°„ë‹¨)"):
    for k, v in GLOSSARY.items():
        st.markdown(f"- **{k}**: {v}")

st.caption("Â© 2025 â€” CRC Agreement Analyzer â€” MAX")
